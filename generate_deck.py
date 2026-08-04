import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Theme Colors based on ui_theme.py (Dark Palette)
COLOR_BG_DARK = RGBColor(0x0D, 0x0E, 0x12)       # Deep system background
COLOR_CARD_BG = RGBColor(0x1C, 0x1C, 0x1E)       # Inset card container
COLOR_CARD_SEC = RGBColor(0x2C, 0x2C, 0x2E)      # Secondary container
COLOR_BORDER = RGBColor(0x38, 0x38, 0x3A)        # Subtle border
COLOR_TEXT_PRI = RGBColor(0xF2, 0xF2, 0xF7)      # Primary text (off-white)
COLOR_TEXT_SEC = RGBColor(0xAE, 0xAE, 0xC0)      # Secondary text (muted gray)
COLOR_TEXT_MUTED = RGBColor(0x6C, 0x6C, 0x70)    # Muted label text

COLOR_BLUE = RGBColor(0x40, 0x9C, 0xFF)          # High contrast blue
COLOR_RED = RGBColor(0xFF, 0x45, 0x3A)           # Danger / Bubble alert red
COLOR_AMBER = RGBColor(0xFF, 0xD6, 0x0A)         # Warning amber / skew
COLOR_GREEN = RGBColor(0x32, 0xD7, 0x4B)         # Fundamental green

FONT_FAMILY = "Inter"

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # Image paths
    workspace_dir = "/Users/danielsflscientific.com/Downloads/Merrill"
    finra_img_path = os.path.join(workspace_dir, "finra.png")
    vol_img_path = os.path.join(workspace_dir, "impliedvolatilitymetric.png")

    def set_slide_background(slide, color=COLOR_BG_DARK):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, section_num, section_name, audience_tag, conf_score):
        # Header bar container
        header_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.9))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Section Category & Badges
        p_sec = tf.paragraphs[0]
        p_sec.text = f"SECTION {section_num}: {section_name.upper()}   |   AUDIENCE: {audience_tag.upper()}   |   CONFIDENCE SCORE: {conf_score:.2f}"
        p_sec.font.name = FONT_FAMILY
        p_sec.font.size = Pt(10)
        p_sec.font.bold = True
        p_sec.font.color.rgb = COLOR_BLUE

        # Main Title
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.name = FONT_FAMILY
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_PRI

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
        return shape

    def add_metric_badge(slide, left, top, width, height, label, value, value_color=COLOR_BLUE):
        add_card(slide, left, top, width, height, bg_color=COLOR_CARD_SEC, border_color=COLOR_BORDER)
        box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.1), Inches(width - 0.2), Inches(height - 0.2))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0

        p_val = tf.paragraphs[0]
        p_val.text = value
        p_val.font.name = FONT_FAMILY
        p_val.font.size = Pt(20)
        p_val.font.bold = True
        p_val.font.color.rgb = value_color
        p_val.alignment = PP_ALIGN.CENTER

        p_lbl = tf.add_paragraph()
        p_lbl.text = label.upper()
        p_lbl.font.name = FONT_FAMILY
        p_lbl.font.size = Pt(9)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = COLOR_TEXT_SEC
        p_lbl.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 0: Title Slide
    # ==========================================
    slide0 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide0)

    # Large Center Card
    add_card(slide0, 1.0, 1.2, 11.333, 5.1, bg_color=COLOR_CARD_BG, border_color=COLOR_BLUE)
    title_box = slide0.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(3.9))
    tf0 = title_box.text_frame
    tf0.word_wrap = True

    p0 = tf0.paragraphs[0]
    p0.text = "MULTIDIMENSIONAL ECONOMETRIC & QUANTITATIVE DETECTION OF MARKET BUBBLES"
    p0.font.name = FONT_FAMILY
    p0.font.size = Pt(26)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_BLUE
    p0.space_after = Pt(14)

    p1 = tf0.add_paragraph()
    p1.text = "A Structural Analysis of the 2026 Macroeconomic Environment & System Implementation"
    p1.font.name = FONT_FAMILY
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_PRI
    p1.space_after = Pt(24)

    p2 = tf0.add_paragraph()
    p2.text = "Author: PhD Econometrics & Quantitative Risk Advisory Team\nTarget Audience: C-Suite Executives, Quantitative Risk Experts & System Engineers\nValidation Threshold: All retained methodologies exceed 0.87 Confidence Score"
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_TEXT_SEC

    # ==========================================
    # GROUP 1: KEY FINDINGS (Score: 0.94)
    # ==========================================

    # SLIDE 1: C-Level Summary (MBA Audience)
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    add_header(slide1, "Executive Summary: Macroeconomic Fragility at S&P 7,500", 1, "Key Findings", "C-Level / MBA Audience", 0.94)

    add_metric_badge(slide1, 0.6, 1.4, 2.8, 1.1, "S&P 500 Level", "7,500 Peak", COLOR_BLUE)
    add_metric_badge(slide1, 3.7, 1.4, 2.8, 1.1, "Shiller CAPE Ratio", "41.37", COLOR_RED)
    add_metric_badge(slide1, 6.8, 1.4, 2.8, 1.1, "FINRA Margin Debt", "$1.416 Trillion", COLOR_RED)
    add_metric_badge(slide1, 9.9, 1.4, 2.8, 1.1, "Buffett Indicator", "218.1% GDP", COLOR_AMBER)

    add_card(slide1, 0.6, 2.8, 5.9, 4.2)
    box1a = slide1.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(5.5), Inches(3.8))
    tf1a = box1a.text_frame
    tf1a.word_wrap = True
    
    p = tf1a.paragraphs[0]
    p.text = "Core Thesis: Valuation Extreme vs. Fundamental Repricing"
    p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(10)

    p = tf1a.add_paragraph()
    p.text = "• S&P 500 near 7,500 creates an unprecedented diagnostic challenge: distinguishing genuine AI technological supercycles from pure speculative exuberance.\n\n• Shiller CAPE at 41.37 represents the 2nd highest valuation epoch in U.S. history, trailing only the 2000 dot-com peak (44.19).\n\n• Market Capitalization to GDP (Buffett Indicator) reached 218.1%, resting 56.6% above long-term economic trendlines."
    p.font.name = FONT_FAMILY; p.font.size = Pt(13); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide1, 6.8, 2.8, 5.9, 4.2)
    box1b = slide1.shapes.add_textbox(Inches(7.0), Inches(3.0), Inches(5.5), Inches(3.8))
    tf1b = box1b.text_frame
    tf1b.word_wrap = True

    p = tf1b.paragraphs[0]
    p.text = "Systemic Risk Catalysts & Options Divergence"
    p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_RED; p.space_after = Pt(10)

    p = tf1b.add_paragraph()
    p.text = "• Systemic Leverage Exhaustion: Margin debt surged 53.7% YoY to $1.416T, depleting excess borrowing buffer ('Margin Credit Exhaustion') and exposing markets to fire sales.\n\n• Options Tail-Risk Alert: While front-month VIX remains suppressed (15-17), CBOE SKEW breached 145+, signaling institutional panic buying of OTM downside protection.\n\n• AI CapEx Grounding: GPT-adjusted econometric models confirm mega-cap AI spend ($754B) has real TFP backing, but peripheral small-caps are in full bubble territory."
    p.font.name = FONT_FAMILY; p.font.size = Pt(13); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 2: Technical Details 1.1 (Expert Audience)
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "Methodological Confidence Scoring & Exclusion Framework", 1, "Key Findings", "Expert Audience", 0.94)

    add_card(slide2, 0.6, 1.4, 5.9, 5.6, border_color=COLOR_RED)
    box2a = slide2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2))
    tf2a = box2a.text_frame
    tf2a.word_wrap = True

    p = tf2a.paragraphs[0]
    p.text = "EXCLUDED METHODOLOGIES (Score < 0.87)"
    p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_RED; p.space_after = Pt(10)

    p = tf2a.add_paragraph()
    p.text = "1. Prediction Markets (Kalshi / Polymarket) — Score: 0.72\n• Empirical study of 5,000+ contracts reveals severe calibration breakdown near expiration.\n• High Brier Scores (>0.25) in key economic sectors.\n• Susceptible to 'Whale' manipulation & reflexivity (betting signal directly distorts real outcome).\n\n2. Extra Trees ML Classifier — Score: 0.86 (86.1% Accuracy)\n• Failed 0.87 threshold for systemic risk precision.\n• Excluded in favor of deep sequence models (LSTM-RNN & HMM) capable of modeling temporal dependencies."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide2, 6.8, 1.4, 5.9, 5.6, border_color=COLOR_GREEN)
    box2b = slide2.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2))
    tf2b = box2b.text_frame
    tf2b.word_wrap = True

    p = tf2b.paragraphs[0]
    p.text = "RETAINED METHODOLOGIES (Score >= 0.87)"
    p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_GREEN; p.space_after = Pt(10)

    p = tf2b.add_paragraph()
    p.text = "• GPT-Adjusted GSADF / PSY Procedure: 0.96 (Gold standard unit root explosive test)\n• Volatility Term Structure Contango: 0.95 (Options pricing mechanics)\n• FINRA Margin Credit Exhaustion: 0.94 (Out-of-sample R^2 = 35.68% annual return predictor)\n• Advanced Payout-Adjusted CAPE (P-CAPE): 0.92 (Explains 35% variance vs 24% standard CAPE)\n• Topological Data Analysis (TDA) & Wavelets: 0.91 (Morlet scaleogram persistent homology)\n• Buffett Indicator (Market Cap/GDP): 0.90 (Macro anchor)\n• Real Estate Price-to-Income / Rent: 0.89 (Affordability limit)\n• NLP LSTM-RNN / HMM Classifiers: 0.89 (Sequence break models)\n• LPPLS Super-Exponential Model: 0.88 (Singularity crash date)"
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 3: Technical Details 1.2 (Expert Audience)
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "Macroeconomic & Valuation Extremes in the 2026 Environment", 1, "Key Findings", "Expert Audience", 0.94)

    add_card(slide3, 0.6, 1.4, 3.8, 5.6)
    box3a = slide3.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(3.4), Inches(5.2))
    tf3a = box3a.text_frame; tf3a.word_wrap = True
    p = tf3a.paragraphs[0]; p.text = "Shiller CAPE & P-CAPE"; p.font.name = FONT_FAMILY; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(8)
    p = tf3a.add_paragraph()
    p.text = "• CAPE = 41.37 (+10.39% YoY).\n• 2nd highest in US history (dot-com peak 44.19).\n• Top quintile CAPE historically yields 0.9% 10-yr annualized real return.\n• Dividend payout fell from 65% (1988) to 35% (2024).\n• P-CAPE adjusts for retained earnings, explaining 35% variance (vs 24% standard CAPE).\n• Earnings yield compressed to ~3.5%, exhausting multiple expansion runway."
    p.font.name = FONT_FAMILY; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide3, 4.7, 1.4, 3.8, 5.6)
    box3b = slide3.shapes.add_textbox(Inches(4.9), Inches(1.6), Inches(3.4), Inches(5.2))
    tf3b = box3b.text_frame; tf3b.word_wrap = True
    p = tf3b.paragraphs[0]; p.text = "Buffett Indicator"; p.font.name = FONT_FAMILY; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_AMBER; p.space_after = Pt(8)
    p = tf3b.add_paragraph()
    p.text = "• Ratio = 218.1% (Wilshire 5000 / Nominal GDP).\n• 56.6% above long-term historical trendline.\n• Globalized corporate revenue argument fails to justify sheer magnitude of current deviation.\n• Multi-decade future cash flows pulled into current pricing.\n• Macro decoupling from domestic output base."
    p.font.name = FONT_FAMILY; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide3, 8.8, 1.4, 3.9, 5.6)
    box3c = slide3.shapes.add_textbox(Inches(9.0), Inches(1.6), Inches(3.5), Inches(5.2))
    tf3c = box3c.text_frame; tf3c.word_wrap = True
    p = tf3c.paragraphs[0]; p.text = "Real Estate Affordability"; p.font.name = FONT_FAMILY; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_RED; p.space_after = Pt(8)
    p = tf3c.add_paragraph()
    p.text = "• Price-to-Income = 7.11x (vs 1990s baseline 3.2x, 2006 peak 7.0x).\n• Demographic gap: Under-40 real home values +30% (2019-2024), real income +9%.\n• NAR Report: $86k income required for starter home vs $75k median household income.\n• Price-to-Rent disequilibrium: Ownership $3,700-$4,300/mo vs Rent $1,450/mo.\n• High Price-to-Rent negatively predicts future price growth."
    p.font.name = FONT_FAMILY; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 4: Technical Details 1.3 (Expert Audience)
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Systemic Leverage & Market Microstructure Fragility", 1, "Key Findings", "Expert Audience", 0.94)

    add_card(slide4, 0.6, 1.4, 6.2, 5.6)
    box4a = slide4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.2))
    tf4a = box4a.text_frame; tf4a.word_wrap = True
    p = tf4a.paragraphs[0]; p.text = "FINRA Margin Debt & Credit Exhaustion"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_RED; p.space_after = Pt(8)
    p = tf4a.add_paragraph()
    p.text = "• Record Nominal Peak: $1.416 Trillion in May 2026 (+53.7% YoY, +8.53% MoM).\n• Real Debt Expansion: +47.4% inflation-adjusted growth over 12 months.\n• Velocity Mismatch: Since 1997 baseline, real margin debt grew 550% vs 358% equity market real growth.\n• Margin Credit Mechanics: Unused borrowing capacity is exhausted. 1 std dev drop in credit predicts -1.1% lower monthly return (Annual out-of-sample R^2 = 35.68%).\n• Pingcang Line Dynamics: Account leverage near maintenance limits triggers forced broker liquidations, driving non-linear positive feedback crashes (1929 & 2015 precedent)."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # Embed Chart Image finra.png
    if os.path.exists(finra_img_path):
        add_card(slide4, 7.0, 1.4, 5.7, 5.6)
        slide4.shapes.add_picture(finra_img_path, Inches(7.1), Inches(1.5), Inches(5.5), Inches(5.4))

    # SLIDE 5: Implementation Details 1.1 (Developer Audience)
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "System Architecture & Logging Engine (`config.py` & `ingestor.py`)", 1, "Key Findings", "Developer Audience", 0.94)

    add_card(slide5, 0.6, 1.4, 5.9, 5.6)
    box5a = slide5.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2))
    tf5a = box5a.text_frame; tf5a.word_wrap = True
    p = tf5a.paragraphs[0]; p.text = "Core Setup & Data Engineering"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(10)
    p = tf5a.add_paragraph()
    p.text = "1. Logging & Domain Exception Hierarchy (`config.py`)\n• RotatingFileHandler logging strictly to `bubble_detector.log`.\n• Custom Exception Classes: DataFetchError, IndicatorComputationError, ModelTrainingError, ValidationError.\n• Standard 2026 Macro Baseline Constants.\n\n2. Efficient Data Ingestion (`ingestor.py`)\n• DataIngestor fetches historical price data via yfinance with fallback synthetic time series.\n• Polars Schema Downcasting: Forces float32 / int32 memory reduction.\n• Missing Value Imputation: Forward fill (fill_null(strategy='forward')) for price series; cubic spline interpolation for low-frequency macro data (GDP, Margin Debt).\n• Parquet Caching: Local .parquet storage for sub-millisecond I/O."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # Code Box
    add_card(slide5, 6.8, 1.4, 5.9, 5.6, bg_color=RGBColor(0x15, 0x16, 0x1A), border_color=COLOR_BLUE)
    box5b = slide5.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2))
    tf5b = box5b.text_frame; tf5b.word_wrap = True
    p = tf5b.paragraphs[0]; p.text = "Python Programmatic Implementation"; p.font.name = FONT_FAMILY; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = COLOR_GREEN; p.space_after = Pt(6)
    p = tf5b.add_paragraph()
    p.text = "```python\n# Polars Downcasting & Parquet Storage\nimport polars as pl\n\nclass DataIngestor:\n    def optimize_schema(self, df: pl.DataFrame) -> pl.DataFrame:\n        return df.with_columns([\n            pl.col(pl.FLOAT_DTYPES).cast(pl.Float32),\n            pl.col(pl.INTEGER_DTYPES).cast(pl.Int32)\n        ])\n\n    def impute_macro(self, df: pl.DataFrame) -> pl.DataFrame:\n        return df.with_columns(\n            pl.col('margin_debt').interpolate(method='cubic')\n        )\n```"
    p.font.name = "Courier New"; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    # ==========================================
    # GROUP 2: SUPPORTING EVIDENCE (Score: 0.93)
    # ==========================================

    # SLIDE 6: C-Level Summary (MBA Audience)
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "Supporting Evidence: Valuation, Econometrics & Machine Learning", 2, "Supporting Evidence", "C-Level / MBA Audience", 0.93)

    add_card(slide6, 0.6, 1.4, 3.8, 5.6)
    box6a = slide6.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(3.4), Inches(5.2))
    tf6a = box6a.text_frame; tf6a.word_wrap = True
    p = tf6a.paragraphs[0]; p.text = "P-CAPE Dividend Correction"; p.font.name = FONT_FAMILY; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(8)
    p = tf6a.add_paragraph()
    p.text = "• Corporate shift from dividends to share buybacks distorted classic Shiller CAPE.\n• Payout-Adjusted CAPE (P-CAPE) brings forward retained earnings.\n• P-CAPE explains 35% of 10-year return variance (vs 24% standard CAPE).\n• Confirms valuation remains severely stretched even under growth-adjusted lens."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide6, 4.7, 1.4, 3.8, 5.6)
    box6b = slide6.shapes.add_textbox(Inches(4.9), Inches(1.6), Inches(3.4), Inches(5.2))
    tf6b = box6b.text_frame; tf6b.word_wrap = True
    p = tf6b.paragraphs[0]; p.text = "GPT AI Econometric Filter"; p.font.name = FONT_FAMILY; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_GREEN; p.space_after = Pt(8)
    p = tf6b.add_paragraph()
    p.text = "• Standard GSADF / PSY unit-root test suffers size distortion during $754B AI CapEx boom.\n• Technology shocks create non-centrality parameters that look like explosive bubbles.\n• Two-step fundamental decomposition filters out rational AI repricing, pinpointing speculative excess elsewhere."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide6, 8.8, 1.4, 3.9, 5.6)
    box6c = slide6.shapes.add_textbox(Inches(9.0), Inches(1.6), Inches(3.5), Inches(5.2))
    tf6c = box6c.text_frame; tf6c.word_wrap = True
    p = tf6c.paragraphs[0]; p.text = "Topological & Wavelet Physics"; p.font.name = FONT_FAMILY; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_AMBER; p.space_after = Pt(8)
    p = tf6c.add_paragraph()
    p.text = "• Topological Data Analysis (TDA) embeds price series into high-dimensional geometric space.\n• Morlet wavelets act as a mathematical microscope detecting structural phase changes.\n• LPPLS models identify super-exponential oscillations preceding critical crash singularities."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 7: Technical Details 2.1 (Expert Audience)
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "P-CAPE Mathematics & Margin Credit Exhaustion", 2, "Supporting Evidence", "Expert Audience", 0.93)

    add_card(slide7, 0.6, 1.4, 5.9, 5.6)
    box7a = slide7.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2))
    tf7a = box7a.text_frame; tf7a.word_wrap = True
    p = tf7a.paragraphs[0]; p.text = "Payout-Adjusted CAPE (P-CAPE) Model"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(8)
    p = tf7a.add_paragraph()
    p.text = "• Formula Logic: Standard CAPE = P / E_10. P-CAPE adjusts past earnings by bringing forward retained earnings (1 - Payout Ratio) compounding at the Cyclically Adjusted Earnings Yield (CAEY).\n• Empirical Proof: Out-of-sample explanatory power R^2 increases from 24% to 35% for prospective 10-year real equity returns.\n• Baseline Implication: Real corporate earnings grew at 2.0% p.a. over 125 years. Current equity yield @ 3.5% leaves zero mathematical runway for further multiple expansion."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide7, 6.8, 1.4, 5.9, 5.6)
    box7b = slide7.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2))
    tf7b = box7b.text_frame; tf7b.word_wrap = True
    p = tf7b.paragraphs[0]; p.text = "Margin Credit Exhaustion & Fire Sale Feedback"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_RED; p.space_after = Pt(8)
    p = tf7b.add_paragraph()
    p.text = "• Margin Credit Definition: Unused borrowing capacity collateralized by paper long gains.\n• Predictive Supremacy: Systematically outperforms price-to-book and P/E ratios out-of-sample. A 1 std dev drop in credit predicts -1.1% monthly return (Annual R^2 = 35.68%).\n• Pingcang Line Dynamics: Account leverage near maintenance limits forces broker liquidations. Indiscriminate selling depresses asset prices, triggering cascading margin calls (1929 & 2015 crash mechanics)."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 8: Technical Details 2.2 (Expert Audience)
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "GPT-Adjusted Econometric Bubble Detection (GSADF / PSY)", 2, "Supporting Evidence", "Expert Audience", 0.93)

    add_card(slide8, 0.6, 1.4, 12.133, 5.6)
    box8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    tf8 = box8.text_frame; tf8.word_wrap = True

    p = tf8.paragraphs[0]; p.text = "Phillips, Shi & Yu (PSY 2015) Procedure & General-Purpose Technology Shock Filtering"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(8)

    p = tf8.add_paragraph()
    p.text = "1. Standard GSADF Test Mechanics:\n• Sequentially executes right-tailed forward recursive augmented Dickey-Fuller (ADF) unit root tests with flexible, expanding window widths.\n• Null Hypothesis H0: Asset price follows random walk with drift. Alternative H1: Mildly explosive process (Evans rational collapsing bubble).\n\n2. The GPT Size Distortion Problem in 2026:\n• Hyperscalers' $754B AI CapEx introduces a non-linear, hump-shaped technology adoption shock into the Campbell-Shiller present-value model.\n• Fundamental price becomes locally explosive during GPT adoption, contaminating the limit distribution with a non-centrality parameter.\n• Result: Unadjusted GSADF mistakenly flags rational AI fundamental repricing as an irrational bubble.\n\n3. Fundamental-versus-Speculative Two-Step Decomposition:\n• Step 1: Regress asset price P_t on empirical technology proxies: Total Factor Productivity (TFP), IT investment, patent grants.\n• Step 2: Run GSADF exclusively on the residual price series (P_t - P_t_fundamental).\n• Outcome: Filters out mega-cap AI false positives while exposing true speculative bubbles in non-earning small-cap tech."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 9: Technical Details 2.3 (Expert Audience)
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "Topological Data Analysis (TDA), Wavelets & LPPLS Models", 2, "Supporting Evidence", "Expert Audience", 0.93)

    add_card(slide9, 0.6, 1.4, 5.9, 5.6)
    box9a = slide9.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2))
    tf9a = box9a.text_frame; tf9a.word_wrap = True
    p = tf9a.paragraphs[0]; p.text = "TDA Persistence & Morlet Wavelets"; p.font.name = FONT_FAMILY; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_AMBER; p.space_after = Pt(8)
    p = tf9a.add_paragraph()
    p.text = "• Takens' Delay Coordinate Embedding: Reconstructs 1D return series into high-dimensional geometric point cloud x_t = (x_t, x_{t-tau}, ..., x_{t-(d-1)tau}).\n• Persistent Homology: Measures appearance/disappearance of topological features (connected components k=0, loops k=1, voids k=2).\n• Morlet Wavelet Scaleogram: Automates sliding window size to adaptive 'Goldilocks' resolution.\n• Early Warning Signal: L^p norms of persistence landscapes exhibit abnormal growth spikes prior to 2000 and 2008 crashes."
    p.font.name = FONT_FAMILY; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide9, 6.8, 1.4, 5.9, 5.6)
    box9b = slide9.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2))
    tf9b = box9b.text_frame; tf9b.word_wrap = True
    p = tf9b.paragraphs[0]; p.text = "LPPLS Singularity & Machine Learning"; p.font.name = FONT_FAMILY; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_GREEN; p.space_after = Pt(8)
    p = tf9b.add_paragraph()
    p.text = "• Deterministic LPPLS Formula:\n  ln(P(t)) = A + B(t_c - t)^m + C(t_c - t)^m cos(w ln(t_c - t) + phi)\n• Bounded NLS Optimization: Solved via Trust-Region Reflective algorithm (Constraints: t_c > t, 0.1 < m < 0.9, 6 < w < 13).\n• LPPLS captures accelerating log-periodic oscillations driven by noise trader positive feedback.\n• Deep Sequence ML (LSTM-RNN & HMM): Predicts structural regime breaks using text sentiment, volatility skewness, and macro variables."
    p.font.name = FONT_FAMILY; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 10: Implementation Details 2.1 (Developer Audience)
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "Quantitative Indicator Module Architecture (`features/`)", 2, "Supporting Evidence", "Developer Audience", 0.93)

    add_card(slide10, 0.6, 1.4, 12.133, 5.6)
    box10 = slide10.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    tf10 = box10.text_frame; tf10.word_wrap = True

    p = tf10.paragraphs[0]; p.text = "Modular Quantitative Feature Extraction Pipeline (`features/`)"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(10)

    p = tf10.add_paragraph()
    p.text = "• technicals.py: Computes Moving Averages (MA20/50/200), RSI (14-day), Bollinger Bands (20-day, 2 std dev), and 20-day rolling volatility.\n\n• macro_valuation.py: Calculates Shiller CAPE (41.37), Payout-Adjusted CAPE (P-CAPE), Buffett Indicator (218.1% GDP), and rolling Z-score metrics.\n\n• leverage.py: Computes FINRA Margin Debt YoY growth, debt velocity, and excess debt capacity ('Margin Credit Exhaustion Score').\n\n• econometric.py: Implements PSY procedure (GSADF t-statistic) integrated with GPT fundamental decomposition to filter out false positive bubble signals on $754B AI CapEx.\n\n• topology.py: Executes Takens' delay coordinate embedding, TDA persistence landscape L2 norm, and Morlet wavelet scaleogram complexity score."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 11: Implementation Details 2.2 (Developer Audience)
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_header(slide11, "Machine Learning Engine & Options Volatility Module", 2, "Supporting Evidence", "Developer Audience", 0.93)

    add_card(slide11, 0.6, 1.4, 5.9, 5.6)
    box11a = slide11.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2))
    tf11a = box11a.text_frame; tf11a.word_wrap = True
    p = tf11a.paragraphs[0]; p.text = "ML Predictor (`structural_breaks.py`)"; p.font.name = FONT_FAMILY; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(8)
    p = tf11a.add_paragraph()
    p.text = "• RobustScaler Preprocessing: Subtracts median and scales by Interquartile Range (IQR), protecting model from extreme outlier volatility spikes.\n• Gradient Boosting Classifier: Predicts forward 20-day drawdown risk probabilities.\n• Expanding-Window Walk-Forward CV: Implements TimeSeriesSplit cross-validation, guaranteeing strictly zero look-ahead bias."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide11, 6.8, 1.4, 5.9, 5.6, bg_color=RGBColor(0x15, 0x16, 0x1A), border_color=COLOR_GREEN)
    box11b = slide11.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2))
    tf11b = box11b.text_frame; tf11b.word_wrap = True
    p = tf11b.paragraphs[0]; p.text = "Python Options Vol Module (`options_vol.py`)"; p.font.name = FONT_FAMILY; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = COLOR_GREEN; p.space_after = Pt(6)
    p = tf11b.add_paragraph()
    p.text = "```python\n# Options & Behavioral Tracking Engine\nclass OptionsVolatilityEngine:\n    def compute_contango_slope(self, vix1d, vix3m):\n        return (vix3m - vix1d) / vix3m\n\n    def skew_alert_trigger(self, skew_val):\n        return skew_val > 145.0  # Tail risk flag\n\n    def cross_asset_vol_ratio(self, ovx, vix):\n        return ovx / vix        # Baseline > 3.0\n```"
    p.font.name = "Courier New"; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    # ==========================================
    # GROUP 3: SECTOR SPECIFIC APPLICATION (Score: 0.91)
    # ==========================================

    # SLIDE 12: C-Level Summary (MBA Audience)
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12)
    add_header(slide12, "Sector Application: AI Supercycle vs. Energy Shock Hazard", 3, "Sector Specific Application", "C-Level / MBA Audience", 0.91)

    add_card(slide12, 0.6, 1.4, 5.9, 5.6)
    box12a = slide12.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2))
    tf12a = box12a.text_frame; tf12a.word_wrap = True
    p = tf12a.paragraphs[0]; p.text = "Tech Sector Concentration Risk"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(10)
    p = tf12a.add_paragraph()
    p.text = "• Semiconductor Concentration: Semi weight in S&P 500 reached 18% (up from 3% a decade ago).\n• Wall Street Earnings Dependency: S&P 500 2026 earnings growth (15%) is driven almost entirely by 86% semiconductor earnings surge.\n• AI CapEx Baseline: Hyperscalers committing $754B in 2026 (projected >$900B in 2027).\n• Small-Cap AI Froth: Russell 2000 & Microcap up 20-25% without TFP gains, creating extreme valuation vulnerability."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide12, 6.8, 1.4, 5.9, 5.6)
    box12b = slide12.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2))
    tf12b = box12b.text_frame; tf12b.word_wrap = True
    p = tf12b.paragraphs[0]; p.text = "Energy Sector Exogenous Threat"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_RED; p.space_after = Pt(10)
    p = tf12b.add_paragraph()
    p.text = "• Fundamental Earnings Paradox: Energy forecasted to post strongest Q2 earnings growth (doubling YoY).\n• Cross-Asset Volatility Spike: CBOE Crude Oil Volatility Index (OVX) spiked 35% in a single session.\n• OVX/VIX Ratio @ 3.5x: Commodities pricing severe geopolitical & inflation risk while equities remain complacent.\n• Deleveraging Trigger: Energy inflation spike would force central bank rate hikes, detonating the $1.4T margin debt bubble."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 13: Technical Details 3.1 (Expert Audience)
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13)
    add_header(slide13, "Technology Concentration & Small-Cap Speculative Froth", 3, "Sector Specific Application", "Expert Audience", 0.91)

    add_card(slide13, 0.6, 1.4, 12.133, 5.6)
    box13 = slide13.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    tf13 = box13.text_frame; tf13.word_wrap = True

    p = tf13.paragraphs[0]; p.text = "Semiconductor Dominance & Valuation Discount Rate Sensitivity"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(10)

    p = tf13.add_paragraph()
    p.text = "1. Index Weight Concentration:\n• Semiconductors now command 18% of S&P 500 total market cap.\n• Mega-cap semiconductor 3-month implied volatility has surged to 73% (more than double 2016 baseline), signaling massive embedded option risk.\n\n2. Structural Disconnect in Secondary Tier Tech:\n• Small-cap indices (Russell 2000 & Microcap) surged +20% and +25% in H1 2026.\n• The rally is heavily driven by unprofitable companies piggybacking on the AI narrative.\n• While hyperscalers' CapEx ($754B) is cointegrated with TFP, secondary tech lacks fundamental cash flow backing.\n\n3. Discount Rate Vulnerability:\n• Valuations in non-earning AI stocks rely entirely on long-dated terminal value cash flow assumptions.\n• Any upward shift in discount rates (driven by rate hikes or inflation) causes violent valuation contraction in secondary tech."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 14: Technical Details 3.2 (Expert Audience)
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14)
    add_header(slide14, "Energy Volatility & Cross-Asset Shock Transmission", 3, "Sector Specific Application", "Expert Audience", 0.91)

    add_card(slide14, 0.6, 1.4, 12.133, 5.6)
    box14 = slide14.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    tf14 = box14.text_frame; tf14.word_wrap = True

    p = tf14.paragraphs[0]; p.text = "The OVX / VIX Anomaly & Macro Contagion Mechanics"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_RED; p.space_after = Pt(10)

    p = tf14.add_paragraph()
    p.text = "1. Cross-Asset Volatility Divergence:\n• CBOE Crude Oil Volatility Index (OVX) single-day spike of +35% elevated the OVX-to-VIX ratio to 3.5x.\n• Options markets are pricing acute supply-chain disruption and commodity inflation exclusively into commodities, bypassing broad equity indices.\n\n2. Margin Debt Deleveraging Transmission Channel:\n• Stage 1: Geopolitical shock triggers sustained energy price surge.\n• Stage 2: Energy inflation compresses corporate profit margins across non-tech sectors and elevates CPI.\n• Stage 3: Central banks lose monetary easing flexibility, holding interest rates higher for longer.\n• Stage 4: High rates increase margin borrowing costs, triggering margin calls at the Pingcang Line and initiating a forced $1.416T deleveraging cascade."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 15: Technical Details 3.3 (Expert Audience)
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15)
    add_header(slide15, "Cross-Sector Contagion & Liquidity Transmission Channels", 3, "Sector Specific Application", "Expert Audience", 0.91)

    add_card(slide15, 0.6, 1.4, 2.8, 5.6)
    box15a = slide15.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(2.4), Inches(5.2))
    tf15a = box15a.text_frame; tf15a.word_wrap = True
    p = tf15a.paragraphs[0]; p.text = "STEP 1: SHOCK"; p.font.name = FONT_FAMILY; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_RED; p.space_after = Pt(6)
    p = tf15a.add_paragraph()
    p.text = "Energy / Geopolitical Shock\n• OVX spikes +35%\n• Oil volatility ratio OVX/VIX hits 3.5x\n• Commodity risk bypasses broad equities initially."
    p.font.name = FONT_FAMILY; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide15, 3.7, 1.4, 2.8, 5.6)
    box15b = slide15.shapes.add_textbox(Inches(3.9), Inches(1.6), Inches(2.4), Inches(5.2))
    tf15b = box15b.text_frame; tf15b.word_wrap = True
    p = tf15b.paragraphs[0]; p.text = "STEP 2: MACRO"; p.font.name = FONT_FAMILY; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_AMBER; p.space_after = Pt(6)
    p = tf15b.add_paragraph()
    p.text = "Inflation & Rate Pressure\n• Energy cost pass-through raises CPI.\n• Central banks forced into hawkish stance.\n• Discount rates rise across all asset classes."
    p.font.name = FONT_FAMILY; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide15, 6.8, 1.4, 2.8, 5.6)
    box15c = slide15.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(2.4), Inches(5.2))
    tf15c = box15c.text_frame; tf15c.word_wrap = True
    p = tf15c.paragraphs[0]; p.text = "STEP 3: LEVERAGE"; p.font.name = FONT_FAMILY; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(6)
    p = tf15c.add_paragraph()
    p.text = "Margin Debt Call Trigger\n• Levered accounts hit Pingcang Line.\n• $1.416T margin debt faces broker calls.\n• Excess Margin Credit is fully exhausted."
    p.font.name = FONT_FAMILY; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide15, 9.9, 1.4, 2.8, 5.6)
    box15d = slide15.shapes.add_textbox(Inches(10.1), Inches(1.6), Inches(2.4), Inches(5.2))
    tf15d = box15d.text_frame; tf15d.word_wrap = True
    p = tf15d.paragraphs[0]; p.text = "STEP 4: CRASH"; p.font.name = FONT_FAMILY; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_RED; p.space_after = Pt(6)
    p = tf15d.add_paragraph()
    p.text = "Indiscriminate Fire Sale\n• Broker liquidations hit liquid mega-cap tech.\n• Implied correlation spikes from <8 to >80.\n• Full systemic market crash."
    p.font.name = FONT_FAMILY; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 16: Implementation Details 3.1 (Developer Audience)
    slide16 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide16)
    add_header(slide16, "Sector Volatility & Dashboard Tab Implementation (`dashboard.py`)", 3, "Sector Specific Application", "Developer Audience", 0.91)

    add_card(slide16, 0.6, 1.4, 5.9, 5.6)
    box16a = slide16.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2))
    tf16a = box16a.text_frame; tf16a.word_wrap = True
    p = tf16a.paragraphs[0]; p.text = "Sector Dashboard Tab 5 Design"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(10)
    p = tf16a.add_paragraph()
    p.text = "• Interactive Tab: 'Sector-Specific Health Dashboard' in NiceGUI application.\n• Semiconductor Trackers: Monitors concentration ratio (18% threshold) and 3-month rolling IV (73% alert).\n• Small-Cap TFP Divergence: Real-time monitoring of Russell 2000 GSADF residual vs TFP growth.\n• Energy Volatility Monitor: Ingests OVX tick data and computes live OVX/VIX ratio with color-coded alert badges (>3.0 = Amber, >3.5 = Red)."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide16, 6.8, 1.4, 5.9, 5.6, bg_color=RGBColor(0x15, 0x16, 0x1A), border_color=COLOR_BLUE)
    box16b = slide16.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2))
    tf16b = box16b.text_frame; tf16b.word_wrap = True
    p = tf16b.paragraphs[0]; p.text = "NiceGUI Component Snippet"; p.font.name = FONT_FAMILY; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = COLOR_GREEN; p.space_after = Pt(6)
    p = tf16b.add_paragraph()
    p.text = "```python\n# NiceGUI Sector Tab Implementation\nfrom nicegui import ui\n\ndef render_sector_tab(data):\n    with ui.tab_panel('Sector-Specific Health'):\n        ui.label('Sector Concentration & Cross-Asset Vol')\n        ovx_vix = data['ovx'] / data['vix']\n        if ovx_vix > 3.5:\n            ui.badge('CRUDE OIL VOLATILITY ALERT',\n                     color='red')\n        # Render Plotly sector dispersion chart\n        ui.plotly(fig_sector_health)\n```"
    p.font.name = "Courier New"; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    # ==========================================
    # GROUP 4: IMPLICATIONS FOR SYSTEMIC STABILITY (Score: 0.95)
    # ==========================================

    # SLIDE 17: C-Level Summary (MBA Audience)
    slide17 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide17)
    add_header(slide17, "Implications for Systemic Stability: Fragility Architecture", 4, "Implications for Systemic Stability", "C-Level / MBA Audience", 0.95)

    add_card(slide17, 0.6, 1.4, 3.8, 5.6)
    box17a = slide17.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(3.4), Inches(5.2))
    tf17a = box17a.text_frame; tf17a.word_wrap = True
    p = tf17a.paragraphs[0]; p.text = "The Concentration Trap"; p.font.name = FONT_FAMILY; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(8)
    p = tf17a.add_paragraph()
    p.text = "• Narrowing Market Breadth: S&P 500 index gains concentrated in a shrinking handful of mega-caps.\n• Illusion of Stability: Index volatility (VIX 15-17) appears calm while underlying stock dispersion (DSPX > 46) is at multi-year highs."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide17, 4.7, 1.4, 3.8, 5.6)
    box17b = slide17.shapes.add_textbox(Inches(4.9), Inches(1.6), Inches(3.4), Inches(5.2))
    tf17b = box17b.text_frame; tf17b.word_wrap = True
    p = tf17b.paragraphs[0]; p.text = "Liquidity Illusion"; p.font.name = FONT_FAMILY; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_RED; p.space_after = Pt(8)
    p = tf17b.add_paragraph()
    p.text = "• Record Margin Debt ($1.416T) has exhausted aggregate unused borrowing capacity.\n• Market lacks a marginal buyer buffer.\n• Any forced selling will encounter a liquidity void, causing rapid price gapping."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide17, 8.8, 1.4, 3.9, 5.6)
    box17c = slide17.shapes.add_textbox(Inches(9.0), Inches(1.6), Inches(3.5), Inches(5.2))
    tf17c = box17c.text_frame; tf17c.word_wrap = True
    p = tf17c.paragraphs[0]; p.text = "Institutional Fear"; p.font.name = FONT_FAMILY; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_AMBER; p.space_after = Pt(8)
    p = tf17c.add_paragraph()
    p.text = "• Smart Money Positioning: CBOE SKEW > 145 reveals aggressive institutional buying of OTM puts.\n• Implied Correlation Collapse (COR3M < 8.0): Historically precedes sharp volatility regime shifts."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 18: Technical Details 4.1 (Expert Audience)
    slide18 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide18)
    add_header(slide18, "Mechanics of Market Liquidity Exhaustion", 4, "Implications for Systemic Stability", "Expert Audience", 0.95)

    add_card(slide18, 0.6, 1.4, 12.133, 5.6)
    box18 = slide18.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    tf18 = box18.text_frame; tf18.word_wrap = True

    p = tf18.paragraphs[0]; p.text = "Margin Credit Depletion & Pingcang Maintenance Collateral Limits"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_RED; p.space_after = Pt(10)

    p = tf18.add_paragraph()
    p.text = "1. Margin Credit Capacity Limits:\n• Margin Credit = Total Allowable Borrowing Power minus Active Margin Debt.\n• In May 2026, active debt ($1.416T) has reached maximum regulatory and brokerage limits, reducing unused credit to historical zero.\n• Result: Levered market participants can no longer step in to buy routine pullbacks.\n\n2. Non-Linear Fire Sale Collateral Dynamics:\n• When account equity hits broker maintenance thresholds (Pingcang Line), automated risk management systems issue immediate margin calls.\n• Investors are forced to sell holdings indiscriminately into a thin bid stack.\n• Liquidations depress asset prices further, triggering a self-reinforcing liquidation loop (1929 & 2015 shadow-financed crash mechanics).\n\n3. Structural Vulnerability Matrix:\n• High Valuation (CAPE 41.37) + High Leverage ($1.416T) + Zero Margin Credit = Maximum Systemic Fragility."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 19: Technical Details 4.2 (Expert Audience)
    slide19 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide19)
    add_header(slide19, "Options Tail-Risk Pricing & Implied Correlation Collapse", 4, "Implications for Systemic Stability", "Expert Audience", 0.95)

    add_card(slide19, 0.6, 1.4, 6.2, 5.6)
    box19a = slide19.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.2))
    tf19a = box19a.text_frame; tf19a.word_wrap = True
    p = tf19a.paragraphs[0]; p.text = "SKEW & Dispersion Indicators"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_AMBER; p.space_after = Pt(8)
    p = tf19a.add_paragraph()
    p.text = "• CBOE SKEW Index > 145 (Peak 154): Reflects extreme pricing of left-tail crash protection relative to ATM options.\n• Normal Bull Range: 100-120. Elevated SKEW while VIX spot remains low (15-17) proves institutional money is aggressively hedging.\n• CBOE Dispersion Index (DSPX) > 46 vs Implied Correlation (COR3M) < 8.0: Multi-year record divergence.\n• Structural Interpretation: Individual stocks are moving independently while index is pinned—a classic late-stage terminal bull characteristic preceding violent correlation spikes."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # Embed Chart Image impliedvolatilitymetric.png
    if os.path.exists(vol_img_path):
        add_card(slide19, 7.0, 1.4, 5.7, 5.6)
        slide19.shapes.add_picture(vol_img_path, Inches(7.1), Inches(1.5), Inches(5.5), Inches(5.4))

    # SLIDE 20: Technical Details 4.3 (Expert Audience)
    slide20 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide20)
    add_header(slide20, "Non-Linear Chaos & Phase Transition Forecasting", 4, "Implications for Systemic Stability", "Expert Audience", 0.95)

    add_card(slide20, 0.6, 1.4, 12.133, 5.6)
    box20 = slide20.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    tf20 = box20.text_frame; tf20.word_wrap = True

    p = tf20.paragraphs[0]; p.text = "Phase Space Attractors & LPPLS Critical Singularity Date"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(10)

    p = tf20.add_paragraph()
    p.text = "1. Topological Phase Shift Detection:\n• Dynamic tracking of TDA persistence landscape L^2 norms reveals a structural shift from a normal gaussian return distribution to a chaotic attractor regime.\n• Scaleogram complexity scores from Morlet wavelets show high-frequency volatility clustering across multiple time horizons.\n\n2. LPPLS Critical Singularity Estimation:\n• Super-exponential power-law fits (ln(P(t)) = A + B(t_c - t)^m + ...) predict critical point t_c where market instability reaches maximum.\n• Accelerating log-periodic price oscillations confirm self-reinforcing noise trader imitation.\n\n3. Operational Risk Regime State Machine:\n• Regime 0: Low Vol Contango (Complacency) -> Regime 1: Dispersion & SKEW Spike (Institutional Hedging) -> Regime 2: Energy/Rate Exogenous Trigger -> Regime 3: Margin Call Liquidation (Systemic Crash)."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 21: Implementation Details 4.1 (Developer Audience)
    slide21 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide21)
    add_header(slide21, "Systemic Risk Dashboard Application Architecture (`dashboard.py`)", 4, "Implications for Systemic Stability", "Developer Audience", 0.95)

    add_card(slide21, 0.6, 1.4, 12.133, 5.6)
    box21 = slide21.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    tf21 = box21.text_frame; tf21.word_wrap = True

    p = tf21.paragraphs[0]; p.text = "5-Tab Interactive Plotly & NiceGUI Integration"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(10)

    p = tf21.add_paragraph()
    p.text = "• Header Bar & CTA Banner: High-impact typography (600-800 weight) with real-time systemic risk status badge.\n• Tab 1: Macro Valuation Dashboard — Interactive Plotly time series of Shiller CAPE, P-CAPE, and Buffett Indicator with historical crash overlays.\n• Tab 2: Liquidity & Leverage Dashboard — FINRA Margin Debt YoY growth, velocity, and live Margin Credit Exhaustion Score gauge.\n• Tab 3: Econometric Bubble Dashboard — GSADF t-statistic charts with GPT fundamental decomposition toggle.\n• Tab 4: Sentiment & Volatility Dashboard — VIX contango term structure, SKEW index alert (>145), and DSPX vs COR3M correlation dispersion tracker.\n• Tab 5: Sector-Specific Health Dashboard — Semiconductor concentration risk and energy volatility ratios."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # ==========================================
    # GROUP 5: STRATEGIC RECOMMENDATIONS (Score: 0.96)
    # ==========================================

    # SLIDE 22: C-Level Summary (MBA Audience)
    slide22 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide22)
    add_header(slide22, "Strategic Recommendations: C-Suite Institutional Roadmap", 5, "Strategic Recommendations", "C-Level / MBA Audience", 0.96)

    add_card(slide22, 0.6, 1.4, 5.9, 2.7)
    box22a = slide22.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.4))
    tf22a = box22a.text_frame; tf22a.word_wrap = True
    p = tf22a.paragraphs[0]; p.text = "1. TDA Wavelet Regime Monitoring"; p.font.name = FONT_FAMILY; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(4)
    p = tf22a.add_paragraph()
    p.text = "Deploy Topological Data Analysis with Morlet wavelets to monitor L^p persistence landscape norms. Systematically reduce equity beta prior to crash onset."
    p.font.name = FONT_FAMILY; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide22, 6.8, 1.4, 5.9, 2.7)
    box22b = slide22.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.4))
    tf22b = box22b.text_frame; tf22b.word_wrap = True
    p = tf22b.paragraphs[0]; p.text = "2. Track FINRA Margin Credit Limit"; p.font.name = FONT_FAMILY; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_RED; p.space_after = Pt(4)
    p = tf22b.add_paragraph()
    p.text = "Treat unused margin credit as a hard liquidity constraint. As capacity exhausts, aggressively adjust downside volatility targets upward to prepare for Pingcang line fire sales."
    p.font.name = FONT_FAMILY; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide22, 0.6, 4.3, 5.9, 2.7)
    box22c = slide22.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(5.5), Inches(2.4))
    tf22c = box22c.text_frame; tf22c.word_wrap = True
    p = tf22c.paragraphs[0]; p.text = "3. GPT-Adjusted Tech Allocation"; p.font.name = FONT_FAMILY; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_GREEN; p.space_after = Pt(4)
    p = tf22c.add_paragraph()
    p.text = "Use GPT fundamental decomposition. Allocate strictly to mega-cap tech backed by CapEx & TFP gains; divest from unprofitable small-cap AI speculative froth."
    p.font.name = FONT_FAMILY; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide22, 6.8, 4.3, 5.9, 2.7)
    box22d = slide22.shapes.add_textbox(Inches(7.0), Inches(4.4), Inches(5.5), Inches(2.4))
    tf22d = box22d.text_frame; tf22d.word_wrap = True
    p = tf22d.paragraphs[0]; p.text = "4. Asymmetric Options & Cross-Asset Hedges"; p.font.name = FONT_FAMILY; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_AMBER; p.space_after = Pt(4)
    p = tf22d.add_paragraph()
    p.text = "Exploit VIX contango via short-term long-gamma positions. Bypass expensive SKEW put options by constructing tail risk hedges using crude oil (OVX) and Treasury (MOVE) volatility."
    p.font.name = FONT_FAMILY; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 23: Technical Details 5.1 (Expert Audience)
    slide23 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide23)
    add_header(slide23, "Institutional Quantitative Risk Protocols", 5, "Strategic Recommendations", "Expert Audience", 0.96)

    add_card(slide23, 0.6, 1.4, 12.133, 5.6)
    box23 = slide23.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    tf23 = box23.text_frame; tf23.word_wrap = True

    p = tf23.paragraphs[0]; p.text = "Algorithmic Trigger Rules & Margin Credit Hard Thresholds"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(10)

    p = tf23.add_paragraph()
    p.text = "1. TDA Persistence Landscape Norm Trigger:\n• Operational Rule: Monitor the L^2 norm of the persistent homology landscape daily.\n• Action Threshold: When L^2 norm breaches +2.5 standard deviations above its 90-day moving median, automatically reduce portfolio equity beta by 30% to 50% within 24 hours.\n\n2. FINRA Margin Credit Exhaustion Rule:\n• Operational Rule: Compute aggregate margin credit velocity (Margin Credit YoY%).\n• Action Threshold: When Margin Credit YoY drops below -15% while nominal debt is at historical zenith, mandate dynamic VaR scaling and double cash reserves.\n\n3. Structural Break Walk-Forward Integration:\n• Integrate `structural_breaks.py` Gradient Boosting predictions with expanding-window walk-forward CV to update portfolio drawdown probability limits dynamically."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 24: Technical Details 5.2 (Expert Audience)
    slide24 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide24)
    add_header(slide24, "GPT-Adjusted Cointegration Asset Allocation Strategy", 5, "Strategic Recommendations", "Expert Audience", 0.96)

    add_card(slide24, 0.6, 1.4, 12.133, 5.6)
    box24 = slide24.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    tf24 = box24.text_frame; tf24.word_wrap = True

    p = tf24.paragraphs[0]; p.text = "Fundamental-versus-Speculative Long/Short Portfolio Construction"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_GREEN; p.space_after = Pt(10)

    p = tf24.add_paragraph()
    p.text = "1. Fundamental Cointegration Decomposition:\n• Decompose asset price: P_t = P_t_fundamental + P_t_speculative.\n• Long Leg: Allocate exclusively to mega-cap technology firms whose price trajectory cointegrates with empirical AI CapEx ($754B) and observable TFP growth.\n• Short Leg: Short secondary small-cap tech (Russell 2000 AI story stocks) that exhibit explosive unadjusted GSADF unit-root stats without corresponding TFP gains.\n\n2. Multiple Re-anchoring via P-CAPE:\n• Replace static forward P/E multiples with Payout-Adjusted CAPE (P-CAPE) to accurately evaluate retained earnings compounding in mega-cap technology balance sheets.\n\n3. Asymmetric Return Profile:\n• Protects upside participation in legitimate technological supercycle while insulating portfolio from speculative bubble collapse."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 25: Technical Details 5.3 (Expert Audience)
    slide25 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide25)
    add_header(slide25, "Asymmetric Options & Cross-Asset Volatility Hedges", 5, "Strategic Recommendations", "Expert Audience", 0.96)

    add_card(slide25, 0.6, 1.4, 12.133, 5.6)
    box25 = slide25.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    tf25 = box25.text_frame; tf25.word_wrap = True

    p = tf25.paragraphs[0]; p.text = "Cross-Asset Volatility Proxy Hedging Matrix"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_AMBER; p.space_after = Pt(10)

    p = tf25.add_paragraph()
    p.text = "1. Exploiting Volatility Term Structure Contango:\n• Front-month VIX1D (<10) vs VIX3M/1Y (19-23) contango allows ultra-cheap short-term long gamma positioning via 1-week straddles/strangles to capture sudden dispersion events.\n\n2. Cross-Asset Volatility Hedges (Bypassing Overpriced SKEW):\n• CBOE SKEW > 145 makes standard OTM S&P 500 put protection exorbitantly expensive.\n• Solution: Construct cross-asset tail risk hedges using Crude Oil Volatility (OVX) and Treasury Volatility (MOVE) call options.\n\n3. Payoff Asymmetry & Cost Efficiency:\n• Reduces negative carry (drag) of hedging by 60% while delivering massive upside payouts during energy/interest rate macroeconomic shock events."
    p.font.name = FONT_FAMILY; p.font.size = Pt(12); p.font.color.rgb = COLOR_TEXT_PRI

    # SLIDE 26: Implementation Details 5.1 (Developer Audience)
    slide26 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide26)
    add_header(slide26, "UI/UX Accessibility Engine & Production Deployment (`ui_theme.py`)", 5, "Strategic Recommendations", "Developer Audience", 0.96)

    add_card(slide26, 0.6, 1.4, 5.9, 5.6)
    box26a = slide26.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2))
    tf26a = box26a.text_frame; tf26a.word_wrap = True
    p = tf26a.paragraphs[0]; p.text = "UI Design System & Accessibility"; p.font.name = FONT_FAMILY; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_BLUE; p.space_after = Pt(10)
    p = tf26a.add_paragraph()
    p.text = "1. WCAG 2.2 AA Programmatic Compliance (`ui_theme.py`)\n• calculate_contrast_ratio & is_wcag_aa_compliant enforce >= 4.5:1 text contrast and >= 3.0:1 UI element contrast.\n\n2. Dyslexia-Friendly Typography Stack\n• Font Stack: SF Pro Text / Inter / OpenDyslexic.\n• Letter Spacing: 0.015em | Line Height: 1.5 body, 1.3 heading.\n• Strictly prohibits decorative fonts.\n\n3. 8px Base Grid Rhythm & iOS 13+ Visual Styling\n• Spacing tokens: 4px, 8px, 16px, 24px, 32px, 48px.\n• Card containers: 14px rounded corners with subtle drop shadows.\n• Dynamic CSS switcher: Synced with Plotly template (plotly_white vs plotly_dark)."
    p.font.name = FONT_FAMILY; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    add_card(slide26, 6.8, 1.4, 5.9, 5.6, bg_color=RGBColor(0x15, 0x16, 0x1A), border_color=COLOR_BLUE)
    box26b = slide26.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2))
    tf26b = box26b.text_frame; tf26b.word_wrap = True
    p = tf26b.paragraphs[0]; p.text = "Python Accessibility Checker Snippet"; p.font.name = FONT_FAMILY; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = COLOR_GREEN; p.space_after = Pt(6)
    p = tf26b.add_paragraph()
    p.text = "```python\n# ui_theme.py WCAG 2.2 Checker\ndef calculate_contrast_ratio(hex1, hex2):\n    lum1 = relative_luminance(parse_hex_color(hex1))\n    lum2 = relative_luminance(parse_hex_color(hex2))\n    l1, l2 = max(lum1, lum2), min(lum1, lum2)\n    return (l1 + 0.05) / (l2 + 0.05)\n\ndef is_wcag_aa_compliant(hex_fg, hex_bg,\n                        is_large=False):\n    ratio = calculate_contrast_ratio(hex_fg, hex_bg)\n    return ratio >= (3.0 if is_large else 4.5)\n```"
    p.font.name = "Courier New"; p.font.size = Pt(11); p.font.color.rgb = COLOR_TEXT_PRI

    # Save presentation
    output_pptx_path = os.path.join(workspace_dir, "MarketBubble_Detection_Presentation_2026.pptx")
    prs.save(output_pptx_path)
    print(f"Successfully generated PowerPoint deck at: {output_pptx_path}")

if __name__ == "__main__":
    create_presentation()
